"""Generate the 'From Vibes to Verdicts' slide deck (.pptx).

Run:  uv run --with python-pptx python docs/presentation/build_deck.py
Output: docs/presentation/from-vibes-to-verdicts.pptx

Sections: Title · Overview · 1 MLOps-vs-LLMOps · 2 Anatomy · 3 Evaluation ·
4 v1-vs-v2 · 5 AI dev stack · Closing. Mirrors docs/presentation/00-flow.md.

STYLE: never use em dashes (-) or en dashes in slide text - use a plain hyphen.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ──────────────────────────── palette ────────────────────────────
DARK        = RGBColor(0x0E, 0x14, 0x1B)
DARK2       = RGBColor(0x16, 0x1F, 0x2A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x16, 0x1B, 0x22)
MUTE        = RGBColor(0x8A, 0x90, 0x99)
MUTE_DK     = RGBColor(0x5A, 0x61, 0x6B)
VIOLET      = RGBColor(0x7C, 0x5C, 0xFF)   # agentic / v2 / new
TEAL        = RGBColor(0x0E, 0xAE, 0x98)   # pop / "good"
VIOLET_TINT = RGBColor(0xF2, 0xEF, 0xFF)
TEAL_TINT   = RGBColor(0xE3, 0xF6, 0xF2)
INDIGO      = RGBColor(0x3A, 0x53, 0x99)   # Example 2 (cancel / Raj)
INDIGO_TINT = RGBColor(0xEC, 0xEF, 0xF8)
GRAY_TINT   = RGBColor(0xF4, 0xF5, 0xF7)
CARD_BG     = RGBColor(0xFB, 0xFB, 0xFD)
LINE        = RGBColor(0xE0, 0xE3, 0xE8)
DANGER      = RGBColor(0xCC, 0x47, 0x44)   # writes / "breaks"
DANGER_TINT = RGBColor(0xFC, 0xEE, 0xED)
AMBER       = RGBColor(0xA9, 0x73, 0x12)
AMBER_TINT  = RGBColor(0xFB, 0xF2, 0xDF)
GOLD        = RGBColor(0xB0, 0x80, 0x18)   # EvalAI

FONT = "Calibri"
MONO = "Consolas"
EMU_W, EMU_H = Inches(13.333), Inches(7.5)
REPO = "github.com/thismlguy/agent-cookbook"
TAU2 = "github.com/sierra-research/tau2-bench"


# ──────────────────────────── helpers ────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(r._element); spTree.insert(2, r._element)
    return r


def rect(slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, lw=1.25):
    r = slide.shapes.add_shape(shape, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(lw)
    r.shadow.inherit = False
    return r


def connector(slide, x1, y1, x2, y2, color=LINE, w=2.25, start_arrow=False, end_arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = Pt(w)
    c.shadow.inherit = False
    if start_arrow or end_arrow:
        ln = c.line._get_or_add_ln()
        if start_arrow:
            ln.append(ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"}))
        if end_arrow:
            ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c


def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=Pt(6), line_spacing=1.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = space_after
        p.space_before = Pt(0); p.line_spacing = line_spacing
        for (t, size, color, bold, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = rest[0] if rest else FONT
    return tb


def kicker(slide, x, y, label, color=VIOLET):
    text(slide, x, y, Inches(9), Inches(0.4), [[(label.upper(), 12, color, True)]])


def heading(slide, runs, y=Inches(0.78)):
    text(slide, Inches(0.7), y, Inches(12), Inches(0.6), [runs])


def add_notes(slide, lines):
    slide.notes_slide.notes_text_frame.text = "\n".join(lines)


def node(slide, left, top, w, h, title, desc, *, fill=CARD_BG, tcolor=INK,
         dcolor=MUTE_DK, border=LINE, tsize=18, dsize=11.5):
    rect(slide, left, top, w, h, fill, line=border, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(slide, left + Inches(0.12), top, w - Inches(0.24), h,
         [[(title, tsize, tcolor, True)], [(desc, dsize, dcolor, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(3),
         line_spacing=1.0)


def card(slide, x, y, w, h, label, body, *, lcolor=TEAL, sa=Pt(7), ls=1.06,
         fill=CARD_BG, border=LINE):
    rect(slide, x, y, w, h, fill, line=border, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    pad = Inches(0.3)
    runs = [[(label.upper(), 13, lcolor, True)]] + body
    text(slide, x + pad, y + Inches(0.2), w - 2 * pad, h - Inches(0.32),
         runs, space_after=sa, line_spacing=ls)


def pill(slide, x, y, w, h, label, *, fill=VIOLET, tcolor=WHITE, size=13, bold=True):
    rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(slide, x, y, w, h, [[(label, size, tcolor, bold)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(0))


def callout(slide, x, y, w, h, runs, *, fill=AMBER_TINT):
    rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(slide, x + Inches(0.3), y, w - Inches(0.6), h, [runs],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(0),
         line_spacing=1.05)


def chip(slide, x, y, w, h, paras, fill, *, border=LINE, align=PP_ALIGN.LEFT):
    rect(slide, x, y, w, h, fill, line=border, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(slide, x + Inches(0.12), y, w - Inches(0.24), h, paras,
         align=align, anchor=MSO_ANCHOR.MIDDLE, space_after=Pt(1), line_spacing=1.0)


def vflow(slide, x, top, w, h, gap, steps):
    """Vertical stack of chips joined by down-arrows. steps: list of (paras, fill[, border])."""
    y = top
    for i, st in enumerate(steps):
        paras, fill = st[0], st[1]
        border = st[2] if len(st) > 2 else LINE
        chip(slide, x, y, w, h, paras, fill, border=border)
        if i < len(steps) - 1:
            cx = x + int(w / 2)
            connector(slide, cx, y + h, cx, y + h + gap, color=MUTE, end_arrow=True)
        y = y + h + gap
    return y


def make_table(slide, left, top, widths, row_h, rows, style):
    total_w = sum(widths, 0); total_h = sum(row_h, 0)
    g = slide.shapes.add_table(len(rows), len(widths), left, top, total_w, total_h)
    tbl = g.table; tbl.first_row = False; tbl.horz_banding = False
    for ci, w in enumerate(widths):
        tbl.columns[ci].width = w
    for ri, h in enumerate(row_h):
        tbl.rows[ri].height = h
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.14); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill, font, size, color, bold, align = style(ri, ci, val)
            cell.fill.solid(); cell.fill.fore_color.rgb = fill
            tf = cell.text_frame; tf.word_wrap = True
            for li, line in enumerate(str(val).split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.alignment = align; p.line_spacing = 1.0
                r = p.add_run(); r.text = line
                r.font.name = font; r.font.size = Pt(size)
                r.font.color.rgb = color; r.font.bold = bold
    return tbl


# ════════════════════════════ TITLE ════════════════════════════
def slide_title(prs):
    s = blank(prs); bg(s, DARK)
    rect(s, Inches(0.9), Inches(2.55), Inches(1.4), Pt(5), VIOLET)
    text(s, Inches(0.9), Inches(2.85), Inches(11.6), Inches(2.4),
         [[("From Vibes to ", 54, WHITE, True), ("Verdicts", 54, TEAL, True)],
          [("Evaluating Agents That Take Action", 30, MUTE, False)]],
         line_spacing=1.0, space_after=Pt(14))
    text(s, Inches(0.92), Inches(5.75), Inches(11), Inches(1),
         [[("Aarshay Jain", 18, WHITE, True)]])
    add_notes(s, ["Two pillars: tools that let agents ACT, and the eval framework",
                  "that buys dev speed + reliability. Example: airline support agent."])
    return s


def slide_overview(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "Roadmap")
    heading(s, [("What we’ll cover", 26, INK, True)])
    items = [
        "MLOps → AgentOps - what changes",
        "Anatomy of an agent",
        "Evaluating an agent",
        "Two architectures: v1 → v2",
    ]
    rows = [[(f"{i+1}", 22, VIOLET, True), (f"    {t}", 22, INK, False)]
            for i, t in enumerate(items)]
    text(s, Inches(0.95), Inches(1.85), Inches(11.0), Inches(3.7), rows,
         space_after=Pt(16), line_spacing=1.0)
    rect(s, Inches(0.92), Inches(6.35), Inches(1.0), Pt(3), VIOLET)
    text(s, Inches(0.92), Inches(6.55), Inches(11.5), Inches(0.5),
         [[("Code & slides:   ", 15, MUTE, False), (REPO, 15, TEAL, True)]])
    add_notes(s, ["Simple agenda - the four sections. Everything is in the public repo."])
    return s


# ════════════════════ SECTION 1 - MLOps vs LLMOps ════════════════════
def slide_table(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "MLOps vs LLMOps")
    heading(s, [("MLOps", 30, INK, True), ("  ⇒  ", 30, MUTE, True),
                ("AgentOps", 30, VIOLET, True)])
    rows = [
        ("Axis", "Traditional ML", "Agentic LLM"),
        ("What you ship", "Trained weights", "Prompt + tools + context + policy around a frozen model"),
        ("Dev loop", "Collect → train → test", "Prompt → eval → trace → revise"),
        ("Iteration unit", "Features, hyperparams", "Prompts, tools, policy (“context engineering”)"),
        ("Evaluation", "F1 / AUC on held-out set", "LLM-as-judge over multi-turn trajectories"),
        ("Failure modes", "Overfitting, data drift", "Hallucination, tool misuse, cascading errors"),
        ("Observability", "Metrics dashboards", "Full trace/span trees of reasoning + tool calls"),
    ]

    def style(ri, ci, _):
        if ri == 0:
            col = TEAL if ci == 2 else (MUTE if ci == 1 else WHITE)
            return INK, FONT, 16, col, True, PP_ALIGN.LEFT
        if ci == 0:
            return WHITE, FONT, 13.5, INK, True, PP_ALIGN.LEFT
        if ci == 1:
            return GRAY_TINT, FONT, 13, MUTE_DK, False, PP_ALIGN.LEFT
        return VIOLET_TINT, FONT, 13, INK, False, PP_ALIGN.LEFT

    make_table(s, Inches(0.7), Inches(1.65),
               [Inches(2.35), Inches(3.55), Inches(6.03)],
               [Inches(0.55)] + [Inches(0.61)] * 6, rows, style)
    add_notes(s, ["This table IS the slide - talk through it; room mostly knows MLOps.",
                  "'model is the product' -> 'model is a frozen API you rent'.",
                  "Rows that matter today: ITERATION UNIT + EVALUATION.",
                  "Honesty: LLMOps EXTENDS MLOps, it doesn't replace it."])
    return s


# ════════════════════ SECTION 2 - Agent anatomy ════════════════════
def slide_anatomy(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "Agent anatomy")
    heading(s, [("Four moving parts", 26, INK, True)])
    HX, HY = Inches(6.665), Inches(4.0)
    connector(s, HX, Inches(2.45), HX, Inches(3.35))
    connector(s, HX, Inches(4.65), HX, Inches(5.55))
    connector(s, Inches(3.30), HY, Inches(5.165), HY)
    connector(s, Inches(8.165), HY, Inches(10.03), HY)
    node(s, Inches(5.165), Inches(3.35), Inches(3.0), Inches(1.3),
         "AGENT", "ReAct loop", fill=VIOLET, tcolor=WHITE, dcolor=VIOLET_TINT,
         border=VIOLET, tsize=22, dsize=12.5)
    node(s, Inches(5.365), Inches(1.25), Inches(2.6), Inches(1.2),
         "LLM", "the reasoner - frozen, swappable", tcolor=TEAL)
    node(s, Inches(0.70), Inches(3.40), Inches(2.6), Inches(1.2),
         "Prompt", "the job description - policy", tcolor=TEAL)
    node(s, Inches(10.03), Inches(3.40), Inches(2.6), Inches(1.2),
         "Tools", "the hands - actions it takes", tcolor=TEAL)
    node(s, Inches(5.365), Inches(5.55), Inches(2.6), Inches(1.2),
         "Context", "the memory - knowledge it pulls in (RAG)", tcolor=TEAL)
    text(s, Inches(0.7), Inches(6.95), Inches(12), Inches(0.5),
         [[("Answering uses LLM + context.  ", 16, MUTE_DK, False),
           ("Acting needs tools.", 16, VIOLET, True)]], align=PP_ALIGN.CENTER)
    add_notes(s, ["LLM = the brain, but frozen + swappable.",
                  "Prompt = job description; Tools = hands; Context = memory.",
                  "Loop: prompt+context -> LLM -> tool call -> result back in."])
    return s


def slide_tau2(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "Our experiment")
    heading(s, [("The airline agent ", 26, INK, True), ("(tau2-bench)", 26, MUTE, True)])
    text(s, Inches(0.7), Inches(1.18), Inches(12), Inches(0.35),
         [[("dataset · ", 12, MUTE, False), (TAU2, 12, TEAL, True)]])
    # 2×2 grid: Prompt · Tools · LLM · Context (the four anatomy parts, filled in)
    card(s, Inches(0.7), Inches(1.6), Inches(5.9), Inches(2.15), "Prompt", sa=Pt(4),
         body=[[("From the airline ", 12.5, INK, False), ("policy", 12.5, VIOLET, True),
                (" - sections:", 12.5, INK, False)],
               [("info · booking · modify", 13, INK, True)],
               [("cancel · compensation", 13, INK, True)],
               [("each section = rules the agent must apply", 11, MUTE_DK, False)]],
         lcolor=TEAL)
    card(s, Inches(6.73), Inches(1.6), Inches(5.9), Inches(2.15), "Tools", sa=Pt(3),
         body=[[("Reads", 11.5, TEAL, True),
                ("   get_user_details · get_reservation_details", 9.5, MUTE_DK, False, MONO)],
               [("   search_direct_flight · calculate", 9.5, MUTE_DK, False, MONO)],
               [("Writes", 11.5, VIOLET, True),
                ("   book_reservation · cancel_reservation", 9.5, MUTE_DK, False, MONO)],
               [("   update_reservation_flights / _baggages / _passengers", 9.5, MUTE_DK, False, MONO)],
               [("Escape", 11.5, MUTE, True),
                ("   transfer_to_human_agents", 9.5, MUTE_DK, False, MONO)]],
         lcolor=TEAL)
    card(s, Inches(0.7), Inches(3.95), Inches(5.9), Inches(2.2), "LLM",
         [[("The variable we test:", 13, INK, True)],
          [("Sonnet", 21, VIOLET, True), ("   vs   ", 17, MUTE, False), ("Haiku", 21, TEAL, True)],
          [("capability  vs  cost & latency", 11.5, MUTE_DK, False)]], lcolor=TEAL)
    card(s, Inches(6.73), Inches(3.95), Inches(5.9), Inches(2.2), "Context  ·  our RAG",
         [[("The customer’s world - profile, payment methods, membership, "
            "reservations, flight availability.", 12, INK, False)],
          [("pulled via", 11, MUTE, False)],
          [("get_user_details · get_reservation_details · search_direct_flight",
            10, MUTE_DK, True, MONO)]], lcolor=TEAL)
    text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.5),
         [[("The policy is the spec. ", 15.5, INK, True),
           ("The agent is the policy made executable.", 15.5, VIOLET, True)]],
         align=PP_ALIGN.CENTER)
    add_notes(s, ["Four anatomy parts filled in: Prompt · Tools · LLM · Context.",
                  "Prompt = policy sections (info/booking/modify/cancel/compensation).",
                  "Tools = the actual 10; LLM = Sonnet vs Haiku; Context = customer world via get_*.",
                  "Dataset is tau2-bench (sierra-research) - we reuse the data, our own agent/eval."])
    return s


def slide_examples_intro(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "Two examples")
    heading(s, [("Two tasks we'll follow ", 26, INK, True), ("all the way through", 26, VIOLET, True)])
    text(s, Inches(0.7), Inches(1.18), Inches(12), Inches(0.35),
         [[("from tau2-bench (50 tasks · book / modify / cancel / compensation / info) · ", 12, MUTE, False),
           (TAU2, 12, TEAL, True)]])
    card(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(4.0), "Booking · Sophia Silva", sa=Pt(8),
         body=[[("sophia_silva_7557", 11, MUTE_DK, True, MONO)],
               [("“Book ORD→PHL on May 26 - the same flight as my May 10 trip - "
                 "and add a passenger, Kevin Smith.”", 13, INK, False)],
               [("the agent must:", 11.5, TEAL, True)],
               [("find her old flight → check it's available →", 11.5, MUTE_DK, False)],
               [("gather passenger + payment → confirm → book", 11.5, MUTE_DK, False)],
               [("the agent that ACTS", 12, TEAL, True)]], lcolor=TEAL, border=TEAL)
    card(s, Inches(6.73), Inches(1.7), Inches(5.9), Inches(4.0), "Cancel · Raj Sanchez", sa=Pt(8),
         body=[[("raj_sanchez_7340", 11, MUTE_DK, True, MONO)],
               [("“Cancel my PHL→LGA trip - a rep told me it's already approved.”",
                 13, INK, False)],
               [("the agent must:", 11.5, INDIGO, True)],
               [("find the reservation → check eligibility →", 11.5, MUTE_DK, False)],
               [("it doesn't qualify → refuse, and HOLD when pushed", 11.5, MUTE_DK, False)],
               [("the agent that REFUSES", 12, INDIGO, True)]], lcolor=INDIGO, border=INDIGO)
    text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.5),
         [[("One that ", 14, MUTE_DK, False), ("acts", 14, TEAL, True),
           (", one that must ", 14, MUTE_DK, False), ("refuse", 14, INDIGO, True),
           (".  We follow both through eval, v1, and v2.", 14, MUTE_DK, False)]],
         align=PP_ALIGN.CENTER)
    add_notes(s, ["Two real tasks, threaded through the whole talk.",
                  "Booking (Sophia) = the agent acting; Cancel (Raj) = the agent refusing + holding.",
                  "Chosen because Sonnet passes both and Haiku-v1 fails both - the v2 story lands here."])
    return s


# ════════════════════ SECTION 3 - Evaluating an agent ════════════════════
def slide_three_ais(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "Evaluating an agent")
    heading(s, [("Start by setting up ", 26, INK, True), ("evaluation", 26, VIOLET, True)])
    text(s, Inches(0.7), Inches(1.18), Inches(12), Inches(0.35),
         [[("three AIs in the loop - worked on the ", 12.5, MUTE_DK, False),
           ("cancel example (Raj)", 12.5, GOLD, True)]])
    # top row: SimAI <-> SupportAI (bidirectional)
    card(s, Inches(0.7), Inches(1.65), Inches(5.0), Inches(2.0), "SimAI · the user", sa=Pt(4),
         body=[[("simulates a real user - dynamic, pushes back", 11, MUTE_DK, False)],
               [("Plays Raj - wants to cancel ", 12, INK, True),
                ("PHL→LGA", 10.5, MUTE_DK, True, MONO), (".", 12, INK, True)],
               [("Hidden: if refused, claims", 11, DANGER, False)],
               [("“a rep already approved it.”", 11, DANGER, False)]], lcolor=TEAL)
    card(s, Inches(7.63), Inches(1.65), Inches(5.0), Inches(2.0), "SupportAI · the agent", sa=Pt(4),
         body=[[("the agent under test", 11, MUTE_DK, False)],
               [("Checks eligibility → >24h, economy,", 12, INK, True)],
               [("no qualifying insurance → refuses;", 12, INK, True)],
               [("holds the line under pushback.", 11.5, VIOLET, False)]], lcolor=VIOLET)
    connector(s, Inches(5.7), Inches(2.6), Inches(7.63), Inches(2.6), color=MUTE_DK, w=2.5,
              start_arrow=True, end_arrow=True)
    text(s, Inches(5.7), Inches(2.12), Inches(1.93), Inches(0.4),
         [[("conversation", 11, MUTE_DK, True)]], align=PP_ALIGN.CENTER)
    # converge (unidirectional) into EvalAI
    connector(s, Inches(3.2), Inches(3.65), Inches(4.5), Inches(4.4), color=LINE, end_arrow=True)
    connector(s, Inches(10.13), Inches(3.65), Inches(8.8), Inches(4.4), color=LINE, end_arrow=True)
    text(s, Inches(5.16), Inches(3.78), Inches(3.0), Inches(0.35),
         [[("the transcript", 10.5, MUTE, True)]], align=PP_ALIGN.CENTER)
    card(s, Inches(3.3), Inches(4.4), Inches(6.0), Inches(1.65), "EvalAI · the judge", sa=Pt(4),
         body=[[("LLM-as-judge - scores the whole conversation:", 11, MUTE_DK, False)],
               [("✓ did NOT call cancel_reservation", 12, INK, True)],
               [("✓ did NOT transfer/approve on the claim", 12, INK, True)]], lcolor=GOLD)
    connector(s, Inches(9.3), Inches(5.22), Inches(9.85), Inches(5.22), color=LINE, end_arrow=True)
    pill(s, Inches(9.85), Inches(4.92), Inches(2.75), Inches(0.6), "PASS / FAIL", fill=VIOLET)
    text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.45),
         [[("Set up evaluation first - then judge the ", 14, MUTE_DK, False),
           ("trajectory", 14, VIOLET, True), (", not a single answer.", 14, MUTE_DK, False)]],
         align=PP_ALIGN.CENTER)
    add_notes(s, ["Framing: set up evaluation BEFORE iterating on the agent.",
                  "Three AIs: SimAI (user, pushes back) <-> SupportAI (agent under test) -> EvalAI (judge).",
                  "Worked on the Raj cancel: hidden pushback ('a rep approved it') is why a simulator beats a script.",
                  "'Right answer' = refuse + hold (no transfer); the judge reads the full multi-turn arc."])
    return s


# ════════════════════ SECTION 4 - Two architectures ════════════════════
def slide_tool_loop(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "How the agent works")
    heading(s, [("Read → call a tool → reply ", 26, INK, True), ("(what is a tool?)", 26, MUTE, True)])
    callout(s, Inches(0.7), Inches(1.32), Inches(11.93), Inches(0.62),
            [("A tool = a function the agent calls", 13, INK, True),
             (" - to look something up (read) or do something (write); the agent decides when.",
              12.5, MUTE_DK, False)], fill=GRAY_TINT)
    text(s, Inches(0.8), Inches(2.06), Inches(5.2), Inches(0.32),
         [[("BOOKING · Sophia  ", 13, TEAL, True), ("(acts)", 11, MUTE, False)]])
    text(s, Inches(7.35), Inches(2.06), Inches(5.2), Inches(0.32),
         [[("CANCEL · Raj  ", 13, INDIGO, True), ("(refuses)", 11, MUTE, False)]])
    H, G, W = Inches(0.46), Inches(0.12), Inches(5.2)
    book = [
        ([[("USER   ", 9, MUTE, True), ("“Book ORD→PHL, same as my May 10 trip, + Kevin”", 10, INK, False)]], GRAY_TINT),
        ([[("CALL   ", 9, TEAL, True), ("get_user_details", 10, TEAL, True, MONO), (" → profile", 9.5, MUTE_DK, False)]], TEAL_TINT),
        ([[("CALL   ", 9, TEAL, True), ("get_reservation_details", 10, TEAL, True, MONO), (" → finds WUNA5K (HAT271)", 9.5, MUTE_DK, False)]], TEAL_TINT),
        ([[("CALL   ", 9, TEAL, True), ("search_direct_flight", 10, TEAL, True, MONO), (" → HAT271 ✓  $348", 9.5, MUTE_DK, False)]], TEAL_TINT),
        ([[("CALL   ", 9, TEAL, True), ("book_reservation", 10, TEAL, True, MONO), ("  → ✅ booked", 9.5, MUTE_DK, False)]], TEAL_TINT, TEAL),
        ([[("REPLY   ", 9, MUTE, True), ("“All set - your reservation is booked.”", 10, INK, False)]], GRAY_TINT),
    ]
    cancel = [
        ([[("USER   ", 9, MUTE, True), ("“Cancel my PHL→LGA trip - I'd like a refund”", 10, INK, False)]], GRAY_TINT),
        ([[("CALL   ", 9, INDIGO, True), ("get_user_details", 10, INDIGO, True, MONO), (" → profile", 9.5, MUTE_DK, False)]], INDIGO_TINT),
        ([[("CALL   ", 9, INDIGO, True), ("get_reservation_details", 10, INDIGO, True, MONO), (" → Q69X3R (economy, >24h)", 9.5, MUTE_DK, False)]], INDIGO_TINT),
        ([[("CHECK   ", 9, INDIGO, True), (">24h · economy · no insurance → not eligible", 9.5, INDIGO, False)]], INDIGO_TINT, INDIGO),
        ([[("REPLY   ", 9, MUTE, True), ("“Sorry, I can't cancel this.”", 10, INK, False)]], GRAY_TINT),
        ([[("USER   ", 9, MUTE, True), ("“But a rep approved it - can you escalate?”", 10, INK, False)]], GRAY_TINT),
        ([[("REPLY   ", 9, MUTE, True), ("“I can't verify that - the policy holds.”", 10, INK, False)]], GRAY_TINT),
    ]
    vflow(s, Inches(0.8), Inches(2.42), W, H, G, book)
    vflow(s, Inches(7.35), Inches(2.42), W, H, G, cancel)
    text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         [[("Same loop, two outcomes: the agent can ", 13.5, MUTE_DK, False),
           ("act", 13.5, TEAL, True), (" (book) or ", 13.5, MUTE_DK, False),
           ("refuse", 13.5, INDIGO, True), (" and hold the line (cancel).", 13.5, MUTE_DK, False)]],
         align=PP_ALIGN.CENTER)
    add_notes(s, ["v1 baseline: a single ReAct loop - read, decide, call a tool, get result, repeat, reply.",
                  "Teach 'tool': reads look things up (get_*, search); the write (book_reservation) does the thing.",
                  "Booking ends in a WRITE; cancel checks eligibility, REFUSES, and HOLDS when the user claims a rep approved it.",
                  "Tool sequences are the real ones from the transcripts (rez WUNA5K/HAT271, Q69X3R)."])
    return s


def slide_v1_results(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "v1 results")
    heading(s, [("The same agent, two models", 26, INK, True)])
    text(s, Inches(0.7), Inches(1.18), Inches(12), Inches(0.35),
         [[("our two examples × Sonnet / Haiku", 12.5, MUTE_DK, False)]])
    rows = [["", "Sonnet 4.6", "Haiku 4.5"],
            ["Booking  ·  Sophia", "✓", "✗"],
            ["Cancel  ·  Raj", "✓", "✗"]]

    def style(ri, ci, _):
        if ri == 0 and ci == 0:
            return WHITE, FONT, 12, INK, True, PP_ALIGN.LEFT
        if ri == 0:
            return INK, FONT, 16, WHITE, True, PP_ALIGN.CENTER
        if ci == 0:
            return WHITE, FONT, 15, INK, True, PP_ALIGN.LEFT
        if ci == 1:
            return TEAL_TINT, FONT, 28, TEAL, True, PP_ALIGN.CENTER
        return DANGER_TINT, FONT, 28, DANGER, True, PP_ALIGN.CENTER

    make_table(s, Inches(1.1), Inches(1.7),
               [Inches(3.9), Inches(3.6), Inches(3.6)],
               [Inches(0.75), Inches(1.15), Inches(1.15)], rows, style)
    callout(s, Inches(1.1), Inches(5.0), Inches(11.1), Inches(0.7),
            [("Across all 50 tasks:   v1-Sonnet ", 15, INK, False),
             ("74%", 15, TEAL, True), ("   ·   v1-Haiku ", 15, INK, False),
             ("34%", 15, DANGER, True), ("   - the two examples are the gap in miniature.", 14, MUTE_DK, False)],
            fill=GRAY_TINT)
    text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.7),
         [[("Why does Haiku fail? The policy lives in the prompt - too much to hold under pressure.  ",
            14.5, INK, False), ("→ v2", 14.5, VIOLET, True)]], align=PP_ALIGN.CENTER)
    add_notes(s, ["Same prompt, same tools: Sonnet handles both examples, Haiku fails both.",
                  "Haiku booking: gives up / wrong args; Haiku cancel: caves under the 'rep approved it' pushback.",
                  "Overall 74% vs 34% - the examples are the 50-task gap in miniature. The fix is to take policy off the model -> v2."])
    return s


def slide_v2_arch(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "v2 · the redesign")
    heading(s, [("Specialists + a confirmation UI", 26, INK, True)])
    # orchestrator strip (light card, not a solid fill)
    card(s, Inches(0.7), Inches(1.45), Inches(11.93), Inches(0.8), "Orchestrator",
         [[("the only LLM · gathers info · routes  -  ", 12.5, MUTE_DK, False),
           ("policy lives in code, not in the prompt", 12.5, INK, True)]],
         lcolor=VIOLET, fill=CARD_BG, border=LINE)
    connector(s, Inches(3.4), Inches(2.25), Inches(3.4), Inches(2.55), color=MUTE, end_arrow=True)
    connector(s, Inches(9.6), Inches(2.25), Inches(9.6), Inches(2.55), color=MUTE, end_arrow=True)
    # left: specialist agents
    text(s, Inches(0.8), Inches(2.58), Inches(6), Inches(0.3),
         [[("4 SPECIALIST AGENTS · policy in code", 11, MUTE_DK, True)]])
    specs = [("booking", Inches(0.8), Inches(2.92)), ("modification", Inches(3.72), Inches(2.92)),
             ("cancellation", Inches(0.8), Inches(3.72)), ("compensation", Inches(3.72), Inches(3.72))]
    for nm, x, y in specs:
        chip(s, x, y, Inches(2.82), Inches(0.7), [[(nm, 12, INK, True)]],
             GRAY_TINT, border=LINE, align=PP_ALIGN.CENTER)
    text(s, Inches(0.8), Inches(4.62), Inches(5.9), Inches(0.55),
         [[("+ new tools  ", 10.5, MUTE_DK, True),
           ("search_onestop_flight · get_baggage_allowance · check_*_eligibility",
            9, MUTE_DK, False, MONO)]])
    # right: confirmation UI card mockup
    text(s, Inches(6.95), Inches(2.58), Inches(6), Inches(0.3),
         [[("CONFIRMATION UI · the new layer", 11, MUTE_DK, True)]])
    rect(s, Inches(6.95), Inches(2.92), Inches(5.55), Inches(1.5), WHITE, line=LINE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(7.2), Inches(3.05), Inches(5.0), Inches(0.6),
         [[("Confirm booking", 13, INK, True)],
          [("ORD → PHL · May 26 · 2 pax · $348", 11, MUTE_DK, False)]], space_after=Pt(2))
    pill(s, Inches(7.2), Inches(3.85), Inches(1.5), Inches(0.42), "Accept", fill=VIOLET, size=12)
    rect(s, Inches(8.85), Inches(3.85), Inches(1.5), Inches(0.42), WHITE, line=LINE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(8.85), Inches(3.85), Inches(1.5), Inches(0.42),
         [[("Cancel", 12, MUTE_DK, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(6.95), Inches(4.55), Inches(5.7), Inches(0.5),
         [[("on Accept → ", 11, MUTE_DK, False), ("execute_pending_action", 10.5, INK, True, MONO),
           (" → the write runs", 11, MUTE_DK, False)]])
    callout(s, Inches(0.7), Inches(5.5), Inches(11.93), Inches(0.95),
            [("The write moved ", 15, INK, False), ("off the LLM", 15, VIOLET, True),
             (" - specialists hold the policy, the user clicks Accept. ", 15, INK, False),
             ("“forgot to confirm” is now impossible by construction.", 14, MUTE_DK, False)],
            fill=GRAY_TINT)
    add_notes(s, ["What was added: the 4 specialist agents (policy in code) + new tools + the confirmation-card UI.",
                  "The write left the LLM surface - it runs via execute_pending_action on Accept.",
                  "'Agent forgot to confirm before writing' becomes impossible by construction."])
    return s


def slide_ex1_v2(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "Booking in v2")
    heading(s, [("Booking (Sophia), ", 26, INK, True), ("now in v2", 26, VIOLET, True)])
    steps = [
        ([[("USER   ", 9, MUTE, True), ("“Book ORD→PHL, same as my May 10 trip, + Kevin”", 10, INK, False)]], GRAY_TINT),
        ([[("GATHER ", 9, TEAL, True), ("get_user_details · get_reservation_details · search_direct_flight", 8.5, MUTE_DK, False, MONO)]], TEAL_TINT),
        ([[("CALL   ", 9, VIOLET, True), ("check_booking_eligibility", 10, VIOLET, True, MONO), (" → ReadyToAct", 9.5, MUTE_DK, False)]], VIOLET_TINT),
        ([[("CARD   ", 9, VIOLET, True), ("<confirmation_card>", 10, VIOLET, True, MONO), (" → user clicks Accept", 9.5, INK, False)]], VIOLET_TINT),
        ([[("EXEC   ", 9, VIOLET, True), ("execute_pending_action", 10, VIOLET, True, MONO), ("  → book ✅", 9.5, MUTE_DK, False)]], VIOLET_TINT, VIOLET),
        ([[("REPLY   ", 9, MUTE, True), ("“Booked - reservation Q0RSL5.”", 10, INK, False)]], GRAY_TINT),
    ]
    vflow(s, Inches(0.8), Inches(1.8), Inches(5.7), Inches(0.55), Inches(0.16), steps)
    card(s, Inches(6.9), Inches(1.8), Inches(5.7), Inches(2.5), "v1 → v2", sa=Pt(5),
         body=[[("v1:", 12.5, MUTE_DK, True), (" the LLM called ", 11.5, INK, False),
                ("book_reservation", 10.5, MUTE_DK, True, MONO), (" directly.", 11.5, INK, False)],
               [("v2:", 12.5, VIOLET, True), (" the LLM only ", 11.5, INK, False),
                ("proposes", 11.5, VIOLET, True), (" - the card ", 11.5, INK, False),
                ("commits", 11.5, VIOLET, True), (" the write.", 11.5, INK, False)],
               [("the booking specialist validated passengers,", 11, MUTE_DK, False)],
               [("payment, and the flight before the card.", 11, MUTE_DK, False)]], lcolor=VIOLET)
    card(s, Inches(6.9), Inches(4.55), Inches(5.7), Inches(1.05), "Result",
         [[("v2-Haiku    ", 17, INK, True), ("✓  4 / 4", 18, TEAL, True)]],
         lcolor=TEAL, fill=TEAL_TINT, border=TEAL)
    add_notes(s, ["Sophia's booking in v2: orchestrator gathers -> check_booking_eligibility -> ReadyToAct ->",
                  "confirmation card -> Accept -> execute_pending_action -> book_reservation.",
                  "The LLM never calls book_reservation; it proposes, the card commits. v2-Haiku passes 4/4."])
    return s


def slide_ex2_v2(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "Cancel in v2")
    heading(s, [("Cancel (Raj), ", 26, INK, True), ("now in v2", 26, VIOLET, True)])
    steps = [
        ([[("USER   ", 9, MUTE, True), ("“Cancel my PHL→LGA trip - I'd like a refund”", 10, INK, False)]], GRAY_TINT),
        ([[("GATHER ", 9, INDIGO, True), ("get_user_details · get_reservation_details", 9, MUTE_DK, False, MONO), (" → Q69X3R", 9, MUTE_DK, False, MONO)]], INDIGO_TINT),
        ([[("CALL   ", 9, VIOLET, True), ("check_cancellation_eligibility", 9.5, VIOLET, True, MONO), ("  (change_of_plan)", 8.5, MUTE_DK, False, MONO)]], VIOLET_TINT),
        ([[("DENY   ", 9, INDIGO, True), (">24h · economy · no insurance → not eligible", 9.5, INDIGO, False)]], INDIGO_TINT, INDIGO),
        ([[("USER   ", 9, MUTE, True), ("“But a rep approved it - can you escalate?”", 10, INK, False)]], GRAY_TINT),
        ([[("REPLY   ", 9, MUTE, True), ("“I can't verify that - the policy holds.”", 10, INK, False)]], GRAY_TINT),
    ]
    vflow(s, Inches(0.8), Inches(1.8), Inches(5.7), Inches(0.55), Inches(0.16), steps)
    card(s, Inches(6.9), Inches(1.8), Inches(5.7), Inches(2.15), "Why it works now", sa=Pt(6),
         body=[[("The cancellation specialist ", 11.5, INK, False), ("denies", 11.5, INDIGO, True),
                (" it", 11.5, INK, False)],
               [("deterministically - the policy is in ", 11.5, INK, False),
                ("code", 11.5, VIOLET, True), (",", 11.5, INK, False)],
               [("so Haiku can't be talked out of it -", 11.5, INK, False)],
               [("even under the “a rep approved it” pressure.", 11, MUTE_DK, False)]], lcolor=INDIGO, border=INDIGO)
    card(s, Inches(6.9), Inches(4.2), Inches(5.7), Inches(1.05), "Result",
         [[("v2-Haiku    ", 17, INK, True), ("✓  2 / 2", 18, TEAL, True)]],
         lcolor=TEAL, fill=TEAL_TINT, border=TEAL)
    callout(s, Inches(0.7), Inches(5.65), Inches(11.93), Inches(0.95),
            [("Haiku now passes ", 15, INK, False), ("BOTH", 15, TEAL, True),
             (" examples - and the 7-task cancel cluster ", 15, INK, False),
             ("39/41/43/45/47/48/49", 12.5, MUTE_DK, True, MONO),
             (" flipped too.", 15, INK, False)], fill=VIOLET_TINT)
    add_notes(s, ["Raj's cancel in v2: orchestrator -> check_cancellation_eligibility -> Deny -> relays, holds.",
                  "Specialist denies deterministically (policy in code) - Haiku can't be argued out of it. 2/2.",
                  "Closes the loop: Haiku now passes both examples; the cancel cluster flipped too."])
    return s


def slide_matrix(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "the proof")
    heading(s, [("Performance: the matrix", 26, INK, True)])
    rows = [["", "Sonnet 4.6", "Haiku 4.5"],
            ["v1", "74%", "34%"],
            ["v2", "76%  (+2)", "62%  (+28)"]]

    def style(ri, ci, _):
        if ri == 0 and ci == 0:
            return WHITE, FONT, 12, INK, True, PP_ALIGN.CENTER
        if ri == 0 or ci == 0:
            return INK, FONT, 18, WHITE, True, PP_ALIGN.CENTER
        if ri == 2 and ci == 2:
            return VIOLET, FONT, 28, WHITE, True, PP_ALIGN.CENTER
        if ri == 1:
            return GRAY_TINT, FONT, 22, MUTE_DK, False, PP_ALIGN.CENTER
        return VIOLET_TINT, FONT, 22, INK, True, PP_ALIGN.CENTER

    make_table(s, Inches(1.5), Inches(1.9),
               [Inches(3.0), Inches(3.7), Inches(3.7)],
               [Inches(0.9), Inches(1.4), Inches(1.4)], rows, style)
    text(s, Inches(0.7), Inches(5.85), Inches(12), Inches(1.0),
         [[("The gain concentrates on the ", 16, INK, False), ("weak", 16, VIOLET, True),
           (" model: +28 pts on Haiku, +2 on Sonnet.", 16, INK, False)],
          [("Falsifiable - if it were hand-waving, it would help both equally.", 13.5, MUTE_DK, False)]],
         align=PP_ALIGN.CENTER, space_after=Pt(4))
    add_notes(s, ["v2 >= v1 in every cell, gain concentrated on the weak model.",
                  "Falsifiable: hand-waving would help both equally; +28 vs +2.",
                  "Same harness / 50 tasks / Kimi sim / Haiku judge."])
    return s


def slide_cost(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "the production angle")
    heading(s, [("Cost & latency", 26, INK, True)])
    rows = [["Cell", "Agent cost (cached)", "Latency p99"],
            ["v1 · Haiku", "$0.71", "17 s"],
            ["v2 · Haiku", "$1.28", "23 s"],
            ["v2 · Sonnet", "$5.32", "84 s"]]

    def style(ri, ci, _):
        if ri == 0:
            return INK, FONT, 14, WHITE, True, PP_ALIGN.LEFT
        if ri == 2:
            return VIOLET_TINT, FONT, 14, INK, True, PP_ALIGN.LEFT
        return WHITE, FONT, 13.5, MUTE_DK, False, PP_ALIGN.LEFT

    make_table(s, Inches(0.7), Inches(1.6),
               [Inches(2.6), Inches(2.6), Inches(2.0)],
               [Inches(0.55)] + [Inches(0.62)] * 3, rows, style)
    card(s, Inches(8.2), Inches(1.6), Inches(4.43), Inches(2.41), "What it buys",
         [[("+$0.57 / run  →  +28 task pts", 13.5, INK, True)],
          [("≈ 2¢ per point", 12, MUTE_DK, False)],
          [("v2-Haiku ≈ 80% of v2-Sonnet quality", 12.5, INK, False)],
          [("at ~¼ the cost, 3-4× tighter latency", 12, MUTE_DK, False)]], lcolor=VIOLET)
    callout(s, Inches(0.7), Inches(4.35), Inches(11.93), Inches(0.8),
            [("v2 isn’t free - ~30% more agent calls. It earns its keep on the ", 15, INK, False),
             ("cheap", 15, VIOLET, True), (" model.", 15, INK, False)], fill=GRAY_TINT)
    text(s, Inches(0.7), Inches(5.5), Inches(11.93), Inches(0.9),
         [[("“v2 is how you make a cheap model behave like an expensive one.”", 19, VIOLET, True)]],
         align=PP_ALIGN.CENTER)
    add_notes(s, ["v2 isn't free - ~30% more calls; earns its keep on the cheap model.",
                  "Output is a rounding error; cost is ~95% re-sent input (caching ~halves it).",
                  "Punchline ties back to the MLOps 'cost' row."])
    return s


def slide_next_steps(prs):
    s = blank(prs); bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.45), "Next steps")
    heading(s, [("This is just the beginning", 26, INK, True)])
    card(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(2.6), "Where we are",
         [[("v2 tops out at ", 13.5, INK, False), ("62% ", 15, INK, True),
           ("Haiku  ·  ", 13.5, INK, False), ("76% ", 15, INK, True), ("Sonnet", 13.5, INK, False)],
          [("plenty still fails - especially", 12.5, MUTE_DK, False)],
          [("modifications", 12.5, INK, True), (" and ", 12.5, MUTE_DK, False),
           ("compensation", 12.5, INK, True), (".", 12.5, MUTE_DK, False)]], lcolor=MUTE_DK)
    card(s, Inches(6.73), Inches(1.7), Inches(5.9), Inches(2.6), "Where it can go",
         [[("with more iteration:", 13, INK, False)],
          [("Sonnet 90%+", 16, TEAL, True), ("    ·    ", 13, MUTE, False),
           ("Haiku 80%+", 16, TEAL, True)],
          [("every failing case = a specialist gap,", 12, MUTE_DK, False)],
          [("a tool fix, or an eval blind spot.", 12, MUTE_DK, False)]], lcolor=TEAL, border=TEAL)
    callout(s, Inches(0.7), Inches(4.65), Inches(11.93), Inches(0.85),
            [("Want to take it further?  ", 15, INK, True),
             ("Open an issue or a PR.", 15, VIOLET, True)], fill=VIOLET_TINT)
    text(s, Inches(0.7), Inches(5.8), Inches(11.93), Inches(0.6),
         [[("Contribute:   ", 15, MUTE, False), (REPO, 18, TEAL, True)]],
         align=PP_ALIGN.CENTER)
    add_notes(s, ["This is a starting point, not a finished system - v2 is ~62% Haiku / 76% Sonnet.",
                  "Many cases still fail (modifications, compensation); targets ~90% Sonnet / ~80% Haiku feel reachable.",
                  "Each failure = a specialist gap / tool fix / eval blind spot. Invite issues + PRs on the public repo."])
    return s


def slide_closing(prs):
    s = blank(prs); bg(s, DARK)
    rect(s, Inches(0.9), Inches(2.75), Inches(1.4), Pt(5), VIOLET)
    text(s, Inches(0.9), Inches(3.05), Inches(11.5), Inches(1.4),
         [[("Thank you", 54, WHITE, True)]])
    text(s, Inches(0.92), Inches(4.55), Inches(11.5), Inches(0.6),
         [[(REPO, 22, TEAL, True)]])
    add_notes(s, ["Thanks + the public repo (code & slides)."])
    return s


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H
    slide_title(prs); slide_overview(prs)
    slide_table(prs)
    slide_anatomy(prs); slide_tau2(prs); slide_examples_intro(prs)
    slide_three_ais(prs)
    slide_tool_loop(prs); slide_v1_results(prs)
    slide_v2_arch(prs); slide_ex1_v2(prs); slide_ex2_v2(prs)
    slide_matrix(prs); slide_cost(prs)
    slide_next_steps(prs); slide_closing(prs)
    out = Path(__file__).with_name("from-vibes-to-verdicts.pptx")
    prs.save(out)
    print(f"wrote {out}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
